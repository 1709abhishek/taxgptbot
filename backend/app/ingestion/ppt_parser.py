import io
from typing import List, Dict, Any
from pptx import Presentation

from .base import BaseParser


class PPTParser(BaseParser):
    """Parser for PowerPoint files - extracts text from slides.

    Supports both old .ppt format (using unstructured) and
    new .pptx format (using python-pptx).
    """

    def parse(self, content: bytes, filename: str) -> List[Dict[str, Any]]:
        """Parse PowerPoint content, extracting text from each slide."""
        # Detect actual file format by magic bytes, not extension
        # ZIP files (pptx) start with PK (0x504B)
        # OLE files (old ppt) start with D0 CF 11 E0
        is_zip = content[:2] == b'PK'
        is_ole = content[:4] == b'\xD0\xCF\x11\xE0'

        if is_ole or (filename.lower().endswith('.ppt') and not filename.lower().endswith('.pptx')):
            return self._parse_old_ppt(content, filename)

        return self._parse_pptx(content, filename)

    def _parse_old_ppt(self, content: bytes, filename: str) -> List[Dict[str, Any]]:
        """Parse old .ppt format using unstructured library."""
        import tempfile
        import os
        from unstructured.partition.ppt import partition_ppt

        chunks = []

        # Write to temp file (unstructured needs file path)
        with tempfile.NamedTemporaryFile(suffix='.ppt', delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        try:
            elements = partition_ppt(filename=tmp_path)

            # Collect all text elements
            slide_texts = []
            for element in elements:
                text = str(element).strip()
                if text and len(text) > 10:  # Skip very short elements
                    slide_texts.append(text)

            # Strategy: Create overlapping context windows
            # Each chunk contains 3-5 consecutive elements for better context
            window_size = 4
            step = 2  # Overlap

            if len(slide_texts) <= window_size:
                # Small PPT - one chunk with all content
                if slide_texts:
                    chunks.append(
                        self._create_chunk(
                            content=f"Presentation: {filename}\n\n" + "\n\n".join(slide_texts),
                            filename=filename,
                            page=1,
                            chunk_type="slide",
                            extra_metadata={"total_elements": len(slide_texts)},
                        )
                    )
            else:
                # Create overlapping windows
                chunk_num = 0
                for i in range(0, len(slide_texts), step):
                    window = slide_texts[i:i + window_size]
                    if len(window) >= 2:  # At least 2 elements
                        chunk_num += 1
                        chunks.append(
                            self._create_chunk(
                                content=f"Presentation: {filename} (section {chunk_num}):\n\n" + "\n\n".join(window),
                                filename=filename,
                                page=chunk_num,
                                chunk_type="slide",
                                extra_metadata={"section": chunk_num, "elements": len(window)},
                            )
                        )
        finally:
            os.unlink(tmp_path)

        return chunks

    def _parse_pptx(self, content: bytes, filename: str) -> List[Dict[str, Any]]:
        """Parse new .pptx format using python-pptx."""
        chunks = []

        prs = Presentation(io.BytesIO(content))

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = self._extract_slide_text(slide)

            if slide_text.strip():
                chunks.append(
                    self._create_chunk(
                        content=f"Slide {slide_num} from {filename}:\n\n{slide_text}",
                        filename=filename,
                        page=slide_num,
                        chunk_type="slide",
                        extra_metadata={"slide_number": slide_num},
                    )
                )

            # Extract tables from slide
            table_chunks = self._extract_tables(slide, filename, slide_num)
            chunks.extend(table_chunks)

        return chunks

    def _extract_slide_text(self, slide) -> str:
        """Extract all text from a slide."""
        texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)

        return "\n\n".join(texts)

    def _extract_tables(
        self, slide, filename: str, slide_num: int
    ) -> List[Dict[str, Any]]:
        """Extract tables from slide and convert to Markdown."""
        chunks = []

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                markdown = self._table_to_markdown(table)

                chunks.append(
                    self._create_chunk(
                        content=f"Table from {filename} (slide {slide_num}):\n\n{markdown}",
                        filename=filename,
                        page=slide_num,
                        chunk_type="table",
                    )
                )

        return chunks

    def _table_to_markdown(self, table) -> str:
        """Convert PowerPoint table to Markdown format."""
        rows = []

        for row_idx, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")

            # Add header separator after first row
            if row_idx == 0:
                separator = "| " + " | ".join(["---"] * len(cells)) + " |"
                rows.append(separator)

        return "\n".join(rows)
